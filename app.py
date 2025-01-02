import streamlit as st
from openai import OpenAI
import tempfile
import os
from pypdf import PdfReader
import docx
import pandas as pd
from pptx import Presentation
import openpyxl
import time

# Configuration de la page Streamlit
st.set_page_config(page_title="AI Document Summarizer", layout="wide")

# Configuration d'OpenAI avec la clé API depuis les secrets
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

def extract_text_from_pdf(file):
    """Extrait le texte d'un fichier PDF."""
    reader = PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text() + "\n"
    return text

def extract_text_from_docx(file):
    """Extrait le texte d'un fichier DOCX."""
    doc = docx.Document(file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_excel(file):
    """Extrait le texte d'un fichier Excel."""
    df = pd.read_excel(file)
    text = ""
    # Ajoute les noms des colonnes
    text += "Colonnes : " + ", ".join(df.columns) + "\n\n"
    # Ajoute un résumé des données
    text += f"Nombre de lignes : {len(df)}\n"
    text += f"Résumé des données :\n"
    # Ajoute les premières lignes
    text += df.head().to_string() + "\n"
    return text

def extract_text_from_pptx(file):
    """Extrait le texte d'une présentation PowerPoint."""
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        text += f"\n--- Nouvelle diapositive ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_file_content(uploaded_file):
    """Extrait le contenu du fichier selon son type."""
    if uploaded_file is None:
        return None
        
    file_type = uploaded_file.type
    
    try:
        if file_type == "application/pdf":
            return extract_text_from_pdf(uploaded_file)
        elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_text_from_docx(uploaded_file)
        elif file_type in ["application/vnd.ms-excel", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
            return extract_text_from_excel(uploaded_file)
        elif file_type in ["application/vnd.ms-powerpoint", "application/vnd.openxmlformats-officedocument.presentationml.presentation"]:
            return extract_text_from_pptx(uploaded_file)
        else:
            return uploaded_file.getvalue().decode("utf-8")
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier : {str(e)}")
        return None

def detect_text_language(text):
    """Détecte la langue du texte via GPT."""
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Tu es un expert en détection de langues. Réponds uniquement par le code de langue ISO 639-1 (fr, en, es, etc.)."},
                {"role": "user", "content": f"Quelle est la langue de ce texte ? Réponds uniquement par le code langue.\n\nTexte: {text[:500]}"}
            ],
            temperature=0,
            max_tokens=2
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        st.error(f"Erreur lors de la détection de la langue : {str(e)}")
        return "fr"  # Langue par défaut

def get_summary_from_openai(text, summary_type, max_length, detect_language=True, input_language=None, output_language="fr"):
    """Obtient un résumé via l'API OpenAI."""
    try:
        # Détection de la langue si nécessaire
        if detect_language:
            detected_lang = detect_text_language(text)
            input_language = detected_lang
            if detected_lang != output_language:
                with st.info(f"Langue détectée : {detected_lang}"):
                    st.write("La traduction sera effectuée.")

        # Construction du prompt avec gestion de la langue
        lang_names = {
            "fr": "français", "en": "anglais", "es": "espagnol",
            "de": "allemand", "it": "italien", "pt": "portugais",
            "nl": "néerlandais", "ru": "russe", "zh": "chinois",
            "ja": "japonais"
        }

        lang_instruction = ""
        if input_language != output_language:
            lang_instruction = f"Traduis en {lang_names[output_language]}. "

        prompt_templates = {
            "vulgarized": f"""{lang_instruction}Résume le texte suivant de manière vulgarisée, en utilisant un langage simple 
            et accessible. Longueur approximative : {max_length} mots.
            
            Texte : {text}""",
            
            "technical": f"""{lang_instruction}Fais un résumé technique du texte suivant, en te concentrant sur les aspects 
            techniques et méthodologiques importants. Longueur approximative : {max_length} mots.
            
            Texte : {text}""",
            
            "bullets": f"""{lang_instruction}Résume les points clés du texte suivant sous forme de liste à puces.
            Maximum {max_length} points importants.
            
            Texte : {text}""",
            
            "executive": f"""{lang_instruction}Génère un executive summary du texte suivant, focalisé sur les points 
            stratégiques et les conclusions principales. Longueur approximative : {max_length} mots.
            
            Texte : {text}"""
        }
        
        prompt = prompt_templates[summary_type]
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Tu es un assistant spécialisé dans le résumé de documents."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1000
        )
        
        return response.choices[0].message.content
    except Exception as e:
        st.error(f"Erreur lors de la génération du résumé : {str(e)}")
        return None

def main():
    st.title("📄 AI Document Summarizer")
    
    # Sidebar pour les paramètres
    with st.sidebar:
        st.header("Paramètres")
        
        # Options de langue
        st.subheader("Langues")
        detect_language = st.checkbox("Détecter la langue automatiquement", value=True)
        
        input_language = None
        if not detect_language:
            input_language = st.selectbox(
                "Langue du document",
                ["fr", "en", "es", "de", "it", "pt", "nl", "ru", "zh", "ja"],
                format_func=lambda x: {
                    "fr": "Français",
                    "en": "Anglais",
                    "es": "Espagnol",
                    "de": "Allemand",
                    "it": "Italien",
                    "pt": "Portugais",
                    "nl": "Néerlandais",
                    "ru": "Russe",
                    "zh": "Chinois",
                    "ja": "Japonais"
                }[x]
            )
        
        output_language = st.selectbox(
            "Langue du résumé",
            ["fr", "en", "es", "de", "it", "pt", "nl", "ru", "zh", "ja"],
            format_func=lambda x: {
                "fr": "Français",
                "en": "Anglais",
                "es": "Espagnol",
                "de": "Allemand",
                "it": "Italien",
                "pt": "Portugais",
                "nl": "Néerlandais",
                "ru": "Russe",
                "zh": "Chinois",
                "ja": "Japonais"
            }[x]
        )
        
        st.subheader("Type de résumé")
        summary_type = st.selectbox(
            "Type de résumé",
            ["vulgarized", "technical", "bullets", "executive"],
            format_func=lambda x: {
                "vulgarized": "Vulgarisé",
                "technical": "Technique",
                "bullets": "Points clés",
                "executive": "Executive Summary"
            }[x]
        )
        
        max_length = st.slider(
            "Longueur maximale",
            min_value=100,
            max_value=1000,
            value=300,
            step=50,
            help="Nombre approximatif de mots pour le résumé"
        )

    # Zone principale
    st.header("Source du texte")
    source_type = st.radio(
        "Choisissez la source du texte",
        ["Fichier", "Texte direct"],
        horizontal=True
    )

    text = None

    if source_type == "Fichier":
        uploaded_file = st.file_uploader(
            "Choisissez un fichier",
            type=["txt", "pdf", "docx", "xlsx", "xls", "pptx", "ppt"]
        )
        
        if uploaded_file is not None:
            with st.spinner("Lecture du fichier..."):
                text = get_file_content(uploaded_file)
    else:
        text = st.text_area(
            "Collez votre texte ici",
            height=200,
            help="Vous pouvez directement coller votre texte ici pour le résumer"
        )

    if text:
        col1, col2 = st.columns([1, 1])
        
        with col1:
            st.subheader("Contenu original")
            with st.expander("Voir le contenu", expanded=False):
                st.text_area("", text, height=400)
        
        if st.button("Générer le résumé"):
            with st.spinner("Génération du résumé en cours..."):
                summary = get_summary_from_openai(
                    text,
                    summary_type,
                    max_length,
                    detect_language=detect_language,
                    input_language=input_language,
                    output_language=output_language
                )
                
                if summary:
                    with col2:
                        st.subheader("Résumé")
                        st.markdown(summary)
                        
                        # Bouton de téléchargement du résumé
                        st.download_button(
                            "Télécharger le résumé",
                            summary,
                            "resume.txt",
                            "text/plain"
                        )

if __name__ == "__main__":
    main()
